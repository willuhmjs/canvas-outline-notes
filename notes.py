#!/usr/bin/env python3
"""Canvas assignment + lecture files -> AI study notes -> Outline.

Idempotent via Outline itself: checks whether a document already exists at
Automatic Notes / <class> / Assignments / <assignment> (or Notes / <filename>);
if so, skips entirely (generate-once, never updates). No shared state with
canvas-sync -- Outline's own contents are the only source of truth.

Assignments are filed under the Assignments subfolder; PDF/PPTX lecture files
posted to Canvas course Files are filed under the Notes subfolder.

Deliberately generates a "study scaffold" (outline/strategy/checklist), never a
finished or complete answer, for the third section of the assignment notes -- see
the prompt below for why.

If any of the three external services (Canvas, Outline, chat.cs.odu.edu) reports a
401/403, raises an obvious high-priority "FIX ME" task in the same Academics
calendar canvas-sync writes to (reuses its DAV credentials) rather than just
failing quietly in a job log nobody looks at. Auto-resolves (marked COMPLETED) the
next time a run completes without any auth failure.
"""
import base64
import hashlib
import html
import io
import json
import os
import re
import ssl
import sys
import time
import urllib.error
import urllib.request
import xml.etree.ElementTree as ET
from datetime import datetime, timezone, timedelta
from urllib.parse import urlencode, urljoin
import zoneinfo

import docx  # python-docx
import pymupdf as fitz  # "import fitz" is a deprecated compat shim scheduled for removal
from pptx import Presentation  # python-pptx
try:
    from youtube_transcript_api import YouTubeTranscriptApi
    from youtube_transcript_api._errors import TranscriptsDisabled, NoTranscriptFound
    _YT_AVAILABLE = True
except ImportError:
    _YT_AVAILABLE = False

# The cluster's Traefik proxy intercepts all outbound HTTPS with a self-signed cert.
# Install a global no-verify opener so every urlopen call goes through without failing.
_ssl_ctx = ssl.create_default_context()
_ssl_ctx.check_hostname = False
_ssl_ctx.verify_mode = ssl.CERT_NONE
urllib.request.install_opener(
    urllib.request.build_opener(urllib.request.HTTPSHandler(context=_ssl_ctx))
)

DISPLAY_TZ = zoneinfo.ZoneInfo(os.environ.get("TZ", "America/New_York"))


def format_due(iso_str):
    """Format a Canvas ISO 8601 due date into a human-readable local time string."""
    if not iso_str:
        return "no due date"
    try:
        dt = datetime.fromisoformat(iso_str.replace("Z", "+00:00")).astimezone(DISPLAY_TZ)
        return dt.strftime("%a %b %-d, %Y %-I:%M %p %Z")
    except Exception:
        return iso_str


CANVAS_BASE_URL = os.environ.get("CANVAS_BASE_URL", "https://canvas.odu.edu").rstrip("/")
CANVAS_API_TOKEN = os.environ.get("CANVAS_API_TOKEN", "")

# Fall back to token file written by token_updater (Docker / bare-metal mode)
if not CANVAS_API_TOKEN:
    _token_file = os.environ.get("TOKEN_FILE", "/data/token.json")
    try:
        with open(_token_file) as _f:
            CANVAS_API_TOKEN = json.load(_f).get("token", "")
    except (FileNotFoundError, json.JSONDecodeError, KeyError):
        pass
CHAT_API_BASE_URL = os.environ.get("CHAT_API_BASE_URL", "https://chat.cs.odu.edu/api/v1").rstrip("/")
CHAT_API_KEY = os.environ["CHAT_API_KEY"]
CHAT_MODEL_TEXT = os.environ.get("CHAT_MODEL_TEXT", "gpt-oss-120b")
CHAT_MODEL_VISION = os.environ.get("CHAT_MODEL_VISION", "gemma-4-31b")
OUTLINE_BASE_URL = os.environ.get("OUTLINE_BASE_URL", "https://outline.will.net").rstrip("/")
OUTLINE_API_TOKEN = os.environ["OUTLINE_API_TOKEN"]
OUTLINE_COLLECTION_NAME = os.environ.get("OUTLINE_COLLECTION_NAME", "Automatic Notes")

ASSIGNMENTS_FOLDER = "Assignments"
NOTES_FOLDER = "Notes"

# Assignments are sorted into three time-based buckets under each Assignments folder.
# Current = due within this many days; anything further out is Future.
CURRENT_WINDOW_DAYS = 14
BUCKET_PAST = "Past"
BUCKET_CURRENT = "Current"
BUCKET_FUTURE = "Future"
BUCKET_ORDER = (BUCKET_CURRENT, BUCKET_FUTURE, BUCKET_PAST)  # display order

# For raise_credential_alarm/resolve_credential_alarm only -- same Academics
# calendar and creds canvas-sync already uses. Alarm is skipped (just logged) if
# these aren't set, so this stays optional rather than a hard dependency.
DAV_BASE_URL = os.environ.get("DAV_BASE_URL", "http://davis.dav.svc.cluster.local:9000").rstrip("/")
DAV_USERNAME = os.environ.get("DAV_USERNAME", "")
DAV_PASSWORD = os.environ.get("DAV_PASSWORD", "")
DAV_CALENDAR_DISPLAYNAME = os.environ.get("DAV_CALENDAR_DISPLAYNAME", "Academics")
ALARM_OBJECT_UID = "canvas-outline-notes-credential-alarm"

# File-based state for incremental sync
STATE_FILE = os.environ.get("STATE_FILE", "/state/canvas-notes-state.json")

MAX_FILE_BYTES = 20 * 1024 * 1024
MAX_PDF_PAGES_AS_IMAGES = 5
MAX_ATTACHMENTS_PER_ASSIGNMENT = 10
MAX_IMAGES_PER_REQUEST = 6
# A generous backstop, not a deliberate throttle -- just cheap insurance against a truly
# pathological one-time backlog (e.g. many overlapping current-term courses) blowing past
# activeDeadlineSeconds. Round-robin scheduling (see main()) means every course gets fair
# progress regardless of this number; a normal semester's per-run queue should clear in one
# run well under this.
MAX_NEW_PER_RUN = 200

CANVAS_HEADERS = {"Authorization": f"Bearer {CANVAS_API_TOKEN}"}
# \w+ (not [A-Za-z]+\d+) for the course-code segment -- confirmed live that codes
# with a trailing letter suffix (CYSE425W, ASTP103N) don't match a
# letters-then-digits-only pattern, silently leaving those courses' names uncleaned.
COURSE_NAME_RE = re.compile(r"^\d+_\w+_\d+\s+(.*)$")
FILE_ID_RE = re.compile(r"/files/(\d+)")
TAG_RE = re.compile(r"<[^>]+>")
DAV_NS = {"d": "DAV:", "c": "urn:ietf:params:xml:ns:caldav"}


class AuthFailure(Exception):
    """A 401/403 from one of the three external services -- distinct from other
    errors so it can propagate past the usual per-item try/except and trigger
    raise_credential_alarm instead of just being logged and skipped."""

    def __init__(self, service, detail):
        self.service = service
        self.detail = detail
        super().__init__(f"{service}: {detail}")


def http_json(method, url, headers=None, body=None, timeout=60):
    req_headers = dict(headers or {})
    data = None
    if body is not None:
        data = json.dumps(body).encode("utf-8")
        req_headers.setdefault("Content-Type", "application/json")
    req = urllib.request.Request(url, data=data, method=method, headers=req_headers)
    try:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            raw = resp.read()
            return resp.status, (json.loads(raw) if raw else None)
    except urllib.error.HTTPError as e:
        raw = e.read()
        try:
            parsed = json.loads(raw) if raw else None
        except Exception:
            parsed = raw.decode("utf-8", errors="replace")
        return e.code, parsed


def parse_next_link(link_header):
    if not link_header:
        return None
    for part in link_header.split(","):
        section = [p.strip() for p in part.split(";")]
        if len(section) < 2 or 'rel="next"' not in section[1:]:
            continue
        url_part = section[0]
        if url_part.startswith("<") and url_part.endswith(">"):
            return url_part[1:-1]
    return None


def current_term_courses(courses):
    """Canvas's enrollment_state=active includes courses from concluded past
    terms too (confirmed live: a Fall 2025 and a Spring 2026 course both still
    showed up as "active" alongside the real current Fall 2026 ones) -- and
    start_at/end_at aren't reliable for filtering (the old Fall 2025 course's
    end_at was set a year out, same as current ones). Instead: the current term
    is whichever enrollment_term_id belongs to the most-recently-*started*
    course, which adapts automatically each semester rather than needing a
    hardcoded term id.
    """
    dated = [(c["start_at"], c) for c in courses if c.get("start_at") and c.get("enrollment_term_id")]
    if not dated:
        return courses
    current_term_id = max(dated, key=lambda pair: pair[0])[1]["enrollment_term_id"]
    return [c for c in courses if c.get("enrollment_term_id") == current_term_id]


def canvas_get_all(path, params=None):
    url = CANVAS_BASE_URL + path
    if params:
        url += "?" + urlencode(params, doseq=True)
    results = []
    while url:
        req = urllib.request.Request(url, headers=CANVAS_HEADERS)
        try:
            with urllib.request.urlopen(req, timeout=30) as resp:
                results.extend(json.loads(resp.read()))
                url = parse_next_link(resp.headers.get("Link"))
        except urllib.error.HTTPError as e:
            if e.code in (401, 403):
                raise AuthFailure("Canvas API token", f"HTTP {e.code} from {url}") from e
            raise
    return results


def canvas_get(path):
    req = urllib.request.Request(CANVAS_BASE_URL + path, headers=CANVAS_HEADERS)
    try:
        with urllib.request.urlopen(req, timeout=30) as resp:
            return json.loads(resp.read())
    except urllib.error.HTTPError as e:
        if e.code in (401, 403):
            raise AuthFailure("Canvas API token", f"HTTP {e.code} from {path}") from e
        raise


def clean_course_name(name):
    m = COURSE_NAME_RE.match(name or "")
    return m.group(1).strip() if m else (name or "").strip()


def html_to_text(raw_html):
    if not raw_html:
        return ""
    text = TAG_RE.sub(" ", raw_html)
    text = html.unescape(text)
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n\s*\n+", "\n\n", text)
    return text.strip()


def extract_file_ids(description_html):
    if not description_html:
        return []
    ids, seen = [], set()
    for m in FILE_ID_RE.finditer(description_html):
        fid = m.group(1)
        if fid not in seen:
            seen.add(fid)
            ids.append(fid)
    return ids


def fetch_canvas_file(file_id):
    meta = canvas_get(f"/api/v1/files/{file_id}")
    size = meta.get("size") or 0
    if size and size > MAX_FILE_BYTES:
        return None
    url = meta.get("url")
    if not url:
        return None
    req = urllib.request.Request(url, headers=CANVAS_HEADERS)
    with urllib.request.urlopen(req, timeout=60) as resp:
        content = resp.read(MAX_FILE_BYTES + 1)
    if len(content) > MAX_FILE_BYTES:
        return None
    content_type = meta.get("content-type") or meta.get("content_type") or "application/octet-stream"
    return {"display_name": meta.get("display_name", f"file-{file_id}"), "content_type": content_type, "bytes": content}


def process_pdf(file_bytes):
    """Returns (extracted_text, list_of_data_urls) -- scanned/image-only pages get
    rendered to PNG and treated as images instead of yielding empty text."""
    text_parts, images = [], []
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    image_page_count = 0
    for page in doc:
        text = page.get_text().strip()
        if text:
            text_parts.append(text)
        elif image_page_count < MAX_PDF_PAGES_AS_IMAGES:
            pix = page.get_pixmap(dpi=150)
            png_b64 = base64.b64encode(pix.tobytes("png")).decode()
            images.append(f"data:image/png;base64,{png_b64}")
            image_page_count += 1
    doc.close()
    return "\n\n".join(text_parts), images


def process_docx(file_bytes):
    """Text-only -- no image/diagram extraction, unlike PDFs and PPTX (see
    process_pptx). Adding that would mean parsing docx's raw XML/media parts
    directly (python-docx doesn't expose it); not done here, scope boundary."""
    document = docx.Document(io.BytesIO(file_bytes))
    parts = [p.text for p in document.paragraphs if p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip(" |"):
                parts.append(row_text)
    return "\n".join(parts)


def process_pptx(file_bytes):
    """Text only (slide text frames, tables, and speaker notes) -- unlike PDFs,
    slide images/diagrams aren't rendered and vision-processed. python-pptx can't
    rasterize a slide the way PyMuPDF can a PDF page; doing that would need a
    separate rendering dependency (e.g. LibreOffice headless), not worth it here.
    """
    presentation = Presentation(io.BytesIO(file_bytes))
    slides_out = []
    for i, slide in enumerate(presentation.slides, 1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        lines.append(row_text)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        if lines or notes:
            block = f"[Slide {i}] " + " / ".join(lines)
            if notes:
                block += f" (speaker notes: {notes})"
            slides_out.append(block)
    return "\n".join(slides_out)


PROMPT_TEMPLATE = """You are helping a student understand and prepare for a course assignment. You must NOT write or solve the assignment for them.

CRITICAL, but with one distinction: don't pretend a specific Canvas-provided detail exists when it doesn't -- never invent fake datasets, fake rubric criteria, fake tool names, or fake links as if Canvas gave them to you, that's actively misleading. Separately, for Topic Notes below, you SHOULD draw on your own genuine subject-matter knowledge to write real educational content -- that's not fabrication, it's real information a tutor would know. Just be upfront when you're inferring the likely topic from limited context (course name, module number) rather than an explicit prompt, e.g. "Module 8 in an intro solar-system course typically covers X -- these notes assume that; confirm against your syllabus."

Course: {course}
Assignment: {name}
Due: {due}
Points possible: {points}
Submission type(s): {submission_types}
{rubric_text}
Assignment instructions:
{description}
{attachment_text}
{thin_content_note}
Produce exactly four markdown sections:

## Summary
Plain-language summary of what's actually being asked, grounded only in the information above. If there's little to go on, say that plainly instead of guessing.

## Topic Notes
This is going into a notes app, so this section should be genuinely useful reference material, not just meta-advice: real explanations, definitions, key facts, formulas, or important distinctions on the underlying topic/concepts this assignment covers. Draw on real subject-matter knowledge here. If the exact topic isn't stated, infer the likely one from the course/module/assignment title and say so plainly rather than presenting a guess as certain.

## Resources
A bullet list of concepts/topics worth further study, or -- if a rubric is given -- what its criteria indicate you should focus on. Do not invent specific external resources, datasets, or named tools that aren't mentioned above.

## Study Scaffold
Do NOT write a finished or complete answer. Instead give ONE of these, as bullet points, whichever best fits the assignment type:
- a structural outline/template of what a response should contain (e.g. "a typical response to this prompt covers: X, then Y, then Z")
- for problem-solving/math/code: a step-by-step STRATEGY for approaching it, or one UNWORKED practice problem of similar type/difficulty -- never the solved original problem
- a self-check checklist derived from the instructions/rubric: the specific things a full-credit response needs to address
"""


PRESENTATION_PROMPT_TEMPLATE = """You are helping a student review lecture material from their course.

Course: {course}
Presentation/File: {name}

Extracted content:
{content}

Produce exactly three markdown sections:

## Summary
A concise overview of the main topics and themes covered in this lecture or presentation.

## Key Concepts
Important terms, definitions, formulas, methods, or ideas covered. Be specific and educational -- draw on real subject-matter knowledge for the underlying concepts; if the exact topic isn't fully clear from the content alone, infer from the course/file name and say so plainly rather than presenting a guess as certain.

## Study Notes
Detailed organized notes for later review. Structure by topic or slide section. Include important details, relationships between concepts, and anything that appears emphasized. Write this as genuinely useful reference material, not just a list of what the slides contain.
"""


TEXT_NOTES_PROMPT_TEMPLATE = """You are helping a student review course content.

Course: {course}
Title: {title}
Source type: {source_label}

Content:
{content}

Produce exactly three markdown sections:

## Summary
A concise overview of the main topics covered in this content.

## Key Concepts
Important terms, definitions, ideas, or takeaways. Draw on real subject-matter knowledge to enrich where the content allows; flag when you're supplementing rather than summarising what's here.

## Study Notes
Detailed organised notes for later review. Focus on what's most educationally valuable. Omit boilerplate, navigation text, or anything that clearly isn't course content.
"""


# Signals that a fetched URL is paywalled, login-gated, or otherwise useless.
# Checked case-insensitively against the stripped text body.
PAYWALL_SIGNALS = (
    "sign in to view", "sign in to access", "log in to view", "log in to access",
    "you must be logged in", "login required", "please log in", "please sign in",
    "subscribe to access", "subscription required", "members only", "premium content",
    "access denied", "403 forbidden", "paywall", "create an account to",
    "register to access", "purchase required", "buy now to access",
    "this content is only available", "you need to be logged in",
)

YT_URL_RE = re.compile(
    r'(?:youtube\.com/(?:watch\?(?:.*&)?v=|embed/|shorts/)|youtu\.be/)'
    r'([a-zA-Z0-9_-]{11})'
)


def format_rubric(rubric):
    if not rubric:
        return ""
    lines = ["Grading rubric (from Canvas -- ground the Study Scaffold's checklist in this):"]
    for criterion in rubric:
        desc = (criterion.get("description") or "").strip()
        points = criterion.get("points")
        long_desc = (criterion.get("long_description") or "").strip()
        lines.append(f"- {desc} ({points} pts)" if points is not None else f"- {desc}")
        if long_desc:
            lines.append(f"  {long_desc}")
    return "\n".join(lines)


def build_messages(course, assignment, description_text, extracted_texts, image_data_urls):
    attachment_text = ""
    if extracted_texts:
        attachment_text = "\n\nExtracted attachment text:\n" + "\n\n---\n\n".join(extracted_texts) + "\n"

    rubric_text = format_rubric(assignment.get("rubric"))
    submission_types = ", ".join(assignment.get("submission_types") or []) or "not specified"
    # Common case, confirmed live: LTI/external-tool assignments (publisher platforms like
    # Connect/MasteringX) expose nothing but a title through the Canvas API -- the real
    # content lives entirely in the external tool. Without this, the model confidently
    # fabricated plausible-sounding but entirely made-up specifics instead of saying so.
    is_thin = len(description_text) < 40 and not extracted_texts and not rubric_text
    thin_content_note = (
        "\nNote: Canvas exposes no real instructions or content for this assignment beyond "
        "its title/metadata above -- common for assignments hosted in an external tool or "
        "publisher platform. Say so plainly rather than inventing specifics.\n"
        if is_thin else ""
    )

    prompt = PROMPT_TEMPLATE.format(
        course=course,
        name=assignment.get("name"),
        due=format_due(assignment.get("due_at")),
        points=assignment.get("points_possible", "n/a"),
        submission_types=submission_types,
        rubric_text=(rubric_text + "\n") if rubric_text else "",
        description=description_text or "(no description provided)",
        attachment_text=attachment_text,
        thin_content_note=thin_content_note,
    )
    if image_data_urls:
        content = [{"type": "text", "text": prompt}]
        for data_url in image_data_urls:
            content.append({"type": "image_url", "image_url": {"url": data_url}})
        return content, True
    return prompt, False


def chat_completion(model, content):
    status, resp = http_json(
        "POST", f"{CHAT_API_BASE_URL}/chat/completions",
        headers={"Authorization": f"Bearer {CHAT_API_KEY}"},
        body={"model": model, "messages": [{"role": "user", "content": content}], "max_tokens": 4000},
        timeout=180,
    )
    if status in (401, 403):
        raise AuthFailure("chat.cs.odu.edu API key", f"HTTP {status}: {resp}")
    if status != 200:
        raise RuntimeError(f"chat completion failed: {status} {resp}")
    return resp["choices"][0]["message"]["content"]


def generate_notes(course_name, assignment, existing_id=None):
    description_html = assignment.get("description") or ""
    description_text = html_to_text(description_html)
    file_ids = extract_file_ids(description_html)[:MAX_ATTACHMENTS_PER_ASSIGNMENT]

    # Compute content hash for change detection
    rubric = assignment.get("rubric") or []
    submission_types = assignment.get("submission_types") or []
    content_to_hash = description_html + json.dumps(rubric) + ",".join(submission_types)
    content_hash = hashlib.sha256(content_to_hash.encode()).hexdigest()[:16]

    extracted_texts, image_data_urls = [], []
    for file_id in file_ids:
        try:
            f = fetch_canvas_file(file_id)
        except AuthFailure:
            raise
        except Exception as exc:
            print(f"WARNING: could not fetch file {file_id}: {exc}", file=sys.stderr)
            continue
        if f is None:
            continue
        ctype = f["content_type"]
        display_lower = f["display_name"].lower()
        if ctype.startswith("image/"):
            image_data_urls.append(f"data:{ctype};base64,{base64.b64encode(f['bytes']).decode()}")
        elif ctype == "application/pdf" or display_lower.endswith(".pdf"):
            text, pdf_images = process_pdf(f["bytes"])
            if text:
                extracted_texts.append(f"[{f['display_name']}]\n{text[:5000]}")
            image_data_urls.extend(pdf_images)
        elif ctype.startswith("text/"):
            extracted_texts.append(f"[{f['display_name']}]\n{f['bytes'].decode('utf-8', errors='replace')[:5000]}")
        elif "wordprocessingml" in ctype or display_lower.endswith(".docx"):
            try:
                text = process_docx(f["bytes"])
            except Exception as exc:
                print(f"WARNING: could not parse docx {f['display_name']}: {exc}", file=sys.stderr)
                text = ""
            if text:
                extracted_texts.append(f"[{f['display_name']}]\n{text[:5000]}")
        elif "presentationml" in ctype or display_lower.endswith(".pptx"):
            try:
                text = process_pptx(f["bytes"])
            except Exception as exc:
                print(f"WARNING: could not parse pptx {f['display_name']}: {exc}", file=sys.stderr)
                text = ""
            if text:
                extracted_texts.append(f"[{f['display_name']}]\n{text[:5000]}")
        # legacy binary .doc/.ppt (pre-OOXML) and anything else: not parsed, left as a
        # Canvas-link-only resource -- python-docx/python-pptx only read the OOXML formats

    image_data_urls = image_data_urls[:MAX_IMAGES_PER_REQUEST]
    content, is_vision = build_messages(course_name, assignment, description_text, extracted_texts, image_data_urls)
    model = CHAT_MODEL_VISION if is_vision else CHAT_MODEL_TEXT
    ai_text = chat_completion(model, content)

    # Confirmed live: Outline's markdown importer doesn't support the trailing-
    # two-spaces-then-newline "hard break" convention -- it renders that pattern
    # as a literal backslash-n text string instead of a line break. Plain
    # paragraph breaks (blank line between each) render correctly instead.
    header = (
        f"**Course:** {course_name}\n\n"
        f"**Due:** {format_due(assignment.get('due_at'))}\n\n"
        f"**Points:** {assignment.get('points_possible', 'n/a')}\n\n"
        f"**Hash:** {content_hash}\n\n"
        f"**Canvas link:** [{assignment.get('name')}]({assignment.get('html_url')})\n\n---\n\n"
    )
    doc_text = header + ai_text

    # If updating an existing document, use documents.update instead of returning text
    if existing_id:
        outline_post("documents.update", {"id": existing_id, "text": doc_text, "publish": True})
        return None
    return doc_text


def generate_presentation_notes(course_name, file_name, file_bytes, content_type):
    """Generate study notes from a PDF or PPTX lecture file.
    Returns markdown string, or None if no extractable content."""
    display_lower = file_name.lower()
    extracted_texts = []
    image_data_urls = []

    if content_type == "application/pdf" or display_lower.endswith(".pdf"):
        text, pdf_images = process_pdf(file_bytes)
        if text:
            extracted_texts.append(text[:8000])
        image_data_urls.extend(pdf_images)
    elif "presentationml" in content_type or display_lower.endswith(".pptx"):
        try:
            text = process_pptx(file_bytes)
        except Exception as exc:
            print(f"WARNING: could not parse pptx {file_name}: {exc}", file=sys.stderr)
            text = ""
        if text:
            extracted_texts.append(text[:8000])

    if not extracted_texts and not image_data_urls:
        return None

    content_text = "\n\n---\n\n".join(extracted_texts) if extracted_texts else "(no text extracted)"
    prompt = PRESENTATION_PROMPT_TEMPLATE.format(
        course=course_name,
        name=file_name,
        content=content_text,
    )

    image_data_urls = image_data_urls[:MAX_IMAGES_PER_REQUEST]
    if image_data_urls:
        content = [{"type": "text", "text": prompt}]
        for data_url in image_data_urls:
            content.append({"type": "image_url", "image_url": {"url": data_url}})
        model = CHAT_MODEL_VISION
    else:
        content = prompt
        model = CHAT_MODEL_TEXT

    ai_text = chat_completion(model, content)

    header = (
        f"**Course:** {course_name}\n\n"
        f"**File:** {file_name}\n\n---\n\n"
    )
    return header + ai_text


def get_youtube_transcript(video_id):
    """Returns (transcript_text, None) or (None, skip_reason)."""
    if not _YT_AVAILABLE:
        return None, "youtube-transcript-api not installed"
    try:
        entries = YouTubeTranscriptApi.get_transcript(
            video_id, languages=["en", "en-US", "en-GB", "a.en"]
        )
        text = " ".join(e["text"].strip() for e in entries)
        return text[:10000], None
    except (TranscriptsDisabled, NoTranscriptFound) as exc:
        return None, f"no transcript available: {exc}"
    except Exception as exc:
        return None, f"transcript error: {exc}"


def fetch_url_text(url, timeout=15, max_bytes=300_000):
    """Fetch a URL and extract clean text.
    Returns (text, None) on success or (None, skip_reason) when the content
    is too short, paywalled, non-textual, or otherwise not worth sending to the LLM.
    """
    try:
        req = urllib.request.Request(
            url,
            headers={"User-Agent": "Mozilla/5.0 (compatible; canvas-notes-bot/1.0)"},
        )
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            if resp.status != 200:
                return None, f"HTTP {resp.status}"
            ct = (resp.headers.get("Content-Type") or "").lower()
            if not any(ct.startswith(t) for t in ("text/html", "text/plain", "application/xhtml")):
                return None, f"non-text content-type: {ct.split(';')[0].strip()}"
            raw = resp.read(max_bytes)
            encoding = "utf-8"
            if "charset=" in ct:
                encoding = ct.split("charset=")[-1].strip().split(";")[0].strip()
            text = html_to_text(raw.decode(encoding, errors="replace"))
    except urllib.error.URLError as exc:
        return None, f"fetch error: {exc}"
    except Exception as exc:
        return None, f"error: {exc}"

    if len(text) < 300:
        return None, f"too little text after extraction ({len(text)} chars)"

    text_lower = text.lower()
    for signal in PAYWALL_SIGNALS:
        if signal in text_lower:
            return None, f"paywall/access-gate detected ({signal!r})"

    return text[:8000], None


def generate_text_notes(course_name, title, text, source_label):
    """Generate study notes from pre-extracted text (Canvas page, YouTube transcript, URL)."""
    prompt = TEXT_NOTES_PROMPT_TEMPLATE.format(
        course=course_name,
        title=title,
        source_label=source_label,
        content=text,
    )
    ai_text = chat_completion(CHAT_MODEL_TEXT, prompt)
    header = (
        f"**Course:** {course_name}\n\n"
        f"**Title:** {title}\n\n"
        f"**Type:** {source_label}\n\n---\n\n"
    )
    return header + ai_text


def build_completion_map(courses):
    """{assignment_id: is_done} from Canvas's own submission/grade data.
    "Done" = submitted, graded, or excused. Reuses courses list to avoid re-fetching."""
    completion = {}
    for course in courses:
        course_id = course.get("id")
        if course_id is None:
            continue
        try:
            assignments = canvas_get_all(
                f"/api/v1/courses/{course_id}/assignments",
                {"include[]": "submission", "per_page": 100},
            )
            for a in assignments:
                submission = a.get("submission") or {}
                done = bool(
                    submission.get("submitted_at")
                    or submission.get("excused")
                    or submission.get("grade") is not None
                )
                completion[str(a.get("id"))] = done
        except Exception:
            # Best effort - continue with other courses if one fails
            continue
    return completion


def outline_post(action, body, max_retries=4):
    """Confirmed live: Outline's self-hosted rate limit is real and gets hit within
    a single run once you're doing collections.list + documents.list per course +
    documents.create per class/assignment. Retry with backoff rather than letting
    one 429 crash the whole job (it used to -- caught by an actual run creating 10
    real documents before dying on the 11th call).
    """
    for attempt in range(max_retries + 1):
        status, resp = http_json(
            "POST", f"{OUTLINE_BASE_URL}/api/{action}",
            headers={"Authorization": f"Bearer {OUTLINE_API_TOKEN}"},
            body=body,
        )
        if status == 429 and attempt < max_retries:
            delay = 5 * (attempt + 1)
            print(f"WARNING: Outline rate-limited on {action}, retrying in {delay}s", file=sys.stderr)
            time.sleep(delay)
            continue
        break
    if status in (401, 403):
        raise AuthFailure("Outline API token", f"HTTP {status} from {action}: {resp}")
    if status >= 400:
        raise RuntimeError(f"Outline {action} failed: {status} {resp}")
    return resp["data"]


def read_state():
    """Read state from JSON file. Returns dict or {} if not found."""
    try:
        with open(STATE_FILE, 'r') as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return {}


def write_state(state):
    """Write state to JSON file."""
    # Ensure directory exists
    state_dir = os.path.dirname(STATE_FILE)
    if state_dir:
        os.makedirs(state_dir, exist_ok=True)

    # Write atomically via temp file
    temp_file = STATE_FILE + '.tmp'
    with open(temp_file, 'w') as f:
        json.dump(state, f, indent=2)
    os.replace(temp_file, STATE_FILE)


def find_or_create_collection(name):
    for c in outline_post("collections.list", {"limit": 100}):
        if c.get("name") == name:
            return c["id"]
    created = outline_post("collections.create", {"name": name, "permission": "read"})
    return created["id"]


def list_all_documents(collection_id):
    docs, offset = [], 0
    while True:
        batch = outline_post("documents.list", {"collectionId": collection_id, "limit": 100, "offset": offset})
        if not batch:
            break
        docs.extend(batch)
        if len(batch) < 100:
            break
        offset += 100
    return docs


def find_child_document(all_docs, parent_id, title):
    for d in all_docs:
        if d.get("title") == title and d.get("parentDocumentId") == parent_id:
            return d["id"]
    return None


def create_document(collection_id, parent_document_id, title, text, icon=None):
    body = {"collectionId": collection_id, "title": title, "text": text, "publish": True}
    if parent_document_id:
        body["parentDocumentId"] = parent_document_id
    if icon:
        body["icon"] = icon
    return outline_post("documents.create", body)["id"]


def assignment_bucket(due_at):
    """Classify an assignment's due date into Past / Current / Future."""
    if not due_at:
        return BUCKET_FUTURE
    try:
        due = datetime.fromisoformat(due_at.replace("Z", "+00:00"))
        days = (due.date() - datetime.now(timezone.utc).date()).days
        if days < 0:
            return BUCKET_PAST
        if days <= CURRENT_WINDOW_DAYS:
            return BUCKET_CURRENT
        return BUCKET_FUTURE
    except Exception:
        return BUCKET_FUTURE


def find_or_create_subfolder(all_docs, collection_id, parent_id, title):
    """Find or create a child document used as a section folder."""
    doc_id = find_child_document(all_docs, parent_id, title)
    if doc_id is None:
        doc_id = create_document(collection_id, parent_id, title, f"# {title}\n")
        all_docs.append({"id": doc_id, "title": title, "parentDocumentId": parent_id})
    return doc_id


def ics_escape(value):
    return value.replace("\\", "\\\\").replace("\n", "\\n").replace(",", "\\,").replace(";", "\\;")


def dav_request(method, url, body=None, headers=None):
    req_headers = {"Authorization": "Basic " + base64.b64encode(f"{DAV_USERNAME}:{DAV_PASSWORD}".encode()).decode()}
    if headers:
        req_headers.update(headers)
    data = body.encode("utf-8") if isinstance(body, str) else body
    req = urllib.request.Request(url, data=data, method=method, headers=req_headers)
    try:
        with urllib.request.urlopen(req, timeout=30) as resp:
            return resp.status, resp.headers, resp.read()
    except urllib.error.HTTPError as e:
        return e.code, e.headers, e.read()


def discover_calendar_href():
    def propfind(url, body, depth="0"):
        status, _h, content = dav_request(
            "PROPFIND", url, body=body,
            headers={"Depth": depth, "Content-Type": "application/xml; charset=utf-8"},
        )
        if status != 207:
            raise RuntimeError(f"PROPFIND {url} failed: {status}")
        return ET.fromstring(content)

    def text_of(el, path):
        found = el.find(path, DAV_NS)
        return found.text if found is not None else None

    root = propfind(
        DAV_BASE_URL + "/dav/",
        '<?xml version="1.0" encoding="utf-8"?><d:propfind xmlns:d="DAV:"><d:prop><d:current-user-principal/></d:prop></d:propfind>',
    )
    principal_href = text_of(root, ".//d:response/d:propstat/d:prop/d:current-user-principal/d:href")
    root = propfind(
        urljoin(DAV_BASE_URL + "/", principal_href),
        '<?xml version="1.0" encoding="utf-8"?><d:propfind xmlns:d="DAV:" xmlns:c="urn:ietf:params:xml:ns:caldav">'
        '<d:prop><c:calendar-home-set/></d:prop></d:propfind>',
    )
    home_href = urljoin(DAV_BASE_URL + "/", text_of(root, ".//d:response/d:propstat/d:prop/c:calendar-home-set/d:href"))
    root = propfind(
        home_href,
        '<?xml version="1.0" encoding="utf-8"?><d:propfind xmlns:d="DAV:"><d:prop><d:displayname/></d:prop></d:propfind>',
        depth="1",
    )
    for response in root.findall(".//d:response", DAV_NS):
        if text_of(response, ".//d:prop/d:displayname") == DAV_CALENDAR_DISPLAYNAME:
            return urljoin(DAV_BASE_URL + "/", text_of(response, "d:href"))
    raise RuntimeError(f"calendar {DAV_CALENDAR_DISPLAYNAME!r} not found")


def build_alarm_vtodo(status, summary, description, completed_stamp=None):
    now = datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ")
    today = datetime.now(timezone.utc).strftime("%Y%m%d")
    lines = [
        "BEGIN:VCALENDAR", "VERSION:2.0", "PRODID:-//homelab//canvas-outline-notes//EN",
        "BEGIN:VTODO",
        f"UID:{ALARM_OBJECT_UID}",
        f"DTSTAMP:{now}",
        f"SUMMARY:{ics_escape(summary)}",
        f"DUE;VALUE=DATE:{today}",
        f"STATUS:{status}",
        "PRIORITY:1",
        f"DESCRIPTION:{ics_escape(description)}",
    ]
    if completed_stamp:
        lines += [f"COMPLETED:{completed_stamp}", "PERCENT-COMPLETE:100"]
    lines += ["END:VTODO", "END:VCALENDAR"]
    return "\r\n".join(lines) + "\r\n"


def alarm_object_url():
    calendar_href = discover_calendar_href()
    base = calendar_href if calendar_href.endswith("/") else calendar_href + "/"
    return urljoin(base, f"{ALARM_OBJECT_UID}.ics")


def raise_credential_alarm(service, detail):
    """Best-effort: a failure here just gets logged, never raised further -- an
    alarm that can't be raised must not mask the original AuthFailure."""
    if not DAV_USERNAME or not DAV_PASSWORD:
        print(f"WARNING: no DAV credentials configured, cannot raise calendar alarm for: {service}: {detail}", file=sys.stderr)
        return
    try:
        object_url = alarm_object_url()
        status, headers, _content = dav_request("GET", object_url)
        put_headers = {"Content-Type": "text/calendar; charset=utf-8"}
        if status == 200 and headers.get("ETag"):
            put_headers["If-Match"] = headers["ETag"]
        elif status == 404:
            put_headers["If-None-Match"] = "*"
        body = build_alarm_vtodo(
            "NEEDS-ACTION",
            f"FIX ME: {service} invalid/expired (canvas-outline-notes)",
            f"canvas-outline-notes failed: {detail}. Fix the credential -- this auto-resolves on the next successful run.",
        )
        dav_request("PUT", object_url, body=body, headers=put_headers)
        print(f"raised calendar alarm: {service} invalid/expired", file=sys.stderr)
    except Exception as exc:
        print(f"WARNING: failed to raise calendar alarm: {exc}", file=sys.stderr)


def resolve_credential_alarm():
    """Best-effort, same reasoning as raise_credential_alarm -- never let this
    turn an otherwise-successful run into a failure."""
    if not DAV_USERNAME or not DAV_PASSWORD:
        return
    try:
        object_url = alarm_object_url()
        status, headers, content = dav_request("GET", object_url)
        if status != 200 or b"STATUS:COMPLETED" in content:
            return
        body = build_alarm_vtodo(
            "COMPLETED",
            "FIX ME: credential invalid/expired (canvas-outline-notes)",
            "Resolved automatically -- a run completed successfully.",
            completed_stamp=datetime.now(timezone.utc).strftime("%Y%m%dT%H%M%SZ"),
        )
        put_headers = {"Content-Type": "text/calendar; charset=utf-8"}
        if headers.get("ETag"):
            put_headers["If-Match"] = headers["ETag"]
        dav_request("PUT", object_url, body=body, headers=put_headers)
    except Exception as exc:
        print(f"WARNING: failed to auto-resolve calendar alarm: {exc}", file=sys.stderr)


def main():
    stats = {"created": 0, "moved": 0, "skipped": 0, "errors": 0, "capped": 0, "locked": 0, "regenerated": 0, "marked_submitted": 0}
    state = read_state()
    try:
        collection_id = find_or_create_collection(OUTLINE_COLLECTION_NAME)
        all_docs = list_all_documents(collection_id)

        courses = current_term_courses(canvas_get_all("/api/v1/courses", {"enrollment_state": "active", "per_page": 100}))

        # Build completion map once for all courses
        completion_map = build_completion_map(courses)

        # Backfill icons on existing assignment docs that are missing or outdated.
        # Uses completion_map directly so it works even on legacy docs without a hash.
        # Submitted (Canvas confirmed done) → ✅, pending → 📝.
        _by_id = {d["id"]: d for d in all_docs}
        _has_children = {d.get("parentDocumentId") for d in all_docs if d.get("parentDocumentId")}
        _bucket_names = {BUCKET_CURRENT, BUCKET_FUTURE, BUCKET_PAST}
        # Build title→assignment_id map from completion_map assignments
        _title_to_done = {}
        for course in courses:
            course_id = course.get("id")
            if not course_id:
                continue
            try:
                for a in canvas_get_all(f"/api/v1/courses/{course_id}/assignments", {"per_page": 100}):
                    _title_to_done[a.get("name", "")] = completion_map.get(str(a.get("id")), False)
            except Exception:
                pass
        for doc in all_docs:
            current_icon = doc.get("icon") or ""
            if doc["id"] in _has_children:
                continue
            parent = _by_id.get(doc.get("parentDocumentId", ""), {})
            if parent.get("title") not in _bucket_names:
                continue
            title = doc.get("title", "")
            done = _title_to_done.get(title)
            if done is None:
                continue  # can't determine, leave alone
            want_icon = "✅" if done else "📝"
            if want_icon != current_icon:
                try:
                    outline_post("documents.update", {"id": doc["id"], "icon": want_icon, "publish": True})
                    doc["icon"] = want_icon
                except Exception:
                    pass

        # Gather each course's pending items (assignments + lecture files) first,
        # then interleave round-robin across courses so no single course monopolizes
        # MAX_NEW_PER_RUN. Each item is a tuple:
        #   ("assignment", course_name, parent_doc_id, name, data)
        #   ("file",       course_name, parent_doc_id, name, data)
        pending_by_course = []
        for course in courses:
            course_id = course.get("id")
            if course_id is None:
                continue
            course_name = clean_course_name(course.get("name") or f"Course {course_id}")

            try:
                class_doc_id = find_or_create_subfolder(all_docs, collection_id, None, course_name)
                assignments_folder_id = find_or_create_subfolder(all_docs, collection_id, class_doc_id, ASSIGNMENTS_FOLDER)
                notes_folder_id = find_or_create_subfolder(all_docs, collection_id, class_doc_id, NOTES_FOLDER)
                # Time-based buckets inside Assignments -- created in display order so
                # Outline shows Current first, then Future, then Past.
                bucket_folder_ids = {}
                for bucket in BUCKET_ORDER:
                    bucket_folder_ids[bucket] = find_or_create_subfolder(
                        all_docs, collection_id, assignments_folder_id, bucket
                    )
            except AuthFailure:
                raise
            except Exception as exc:
                print(f"ERROR creating folders for '{course_name}': {exc}", file=sys.stderr)
                stats["errors"] += 1
                continue

            # All parent IDs where assignment docs may live (new buckets + old flat
            # Assignments folder for migration of pre-bucket docs)
            all_assignment_parents = set(bucket_folder_ids.values()) | {assignments_folder_id}

            course_items = []

            # Assignments -> bucketed sub-folders, with re-categorisation on each run
            try:
                # Use incremental sync if we have a last run timestamp
                params = {"per_page": 100}
                last_assignments_run = state.get("last_assignments_run")
                if last_assignments_run:
                    params["updated_since"] = last_assignments_run
                assignments = canvas_get_all(f"/api/v1/courses/{course_id}/assignments", params)
            except AuthFailure:
                raise
            except Exception as exc:
                print(f"ERROR listing assignments for course {course_id}: {exc}", file=sys.stderr)
                stats["errors"] += 1
                assignments = []

            for assignment in assignments:
                name = assignment.get("name") or f"Assignment {assignment.get('id')}"
                target_bucket = assignment_bucket(assignment.get("due_at"))
                target_folder_id = bucket_folder_ids[target_bucket]

                # Find the doc in any bucket (or the old flat Assignments folder)
                existing_id, existing_parent = None, None
                for d in all_docs:
                    if d.get("title") == name and d.get("parentDocumentId") in all_assignment_parents:
                        existing_id = d["id"]
                        existing_parent = d.get("parentDocumentId")
                        break

                if existing_id is not None:
                    if existing_parent != target_folder_id:
                        # Wrong bucket (or pre-migration flat location) -- move it
                        course_items.append(("move", course_name, target_folder_id, name, existing_id))
                    else:
                        # Check if content has changed by comparing hashes
                        try:
                            doc_info = outline_post("documents.info", {"id": existing_id})
                            existing_text = doc_info.get("text", "")
                            hash_match = re.search(r'\*\*Hash:\*\*\s*([a-f0-9]+)', existing_text)
                            if hash_match:
                                existing_hash = hash_match.group(1)
                                # Compute current hash
                                rubric = assignment.get("rubric") or []
                                submission_types = assignment.get("submission_types") or []
                                description_html = assignment.get("description") or ""
                                content_to_hash = description_html + json.dumps(rubric) + ",".join(submission_types)
                                current_hash = hashlib.sha256(content_to_hash.encode()).hexdigest()[:16]
                                if existing_hash != current_hash:
                                    # Content changed, regenerate
                                    course_items.append(("regen", course_name, target_folder_id, name, assignment, existing_id))
                                else:
                                    # Check if should mark as submitted
                                    assignment_id = str(assignment.get("id", ""))
                                    due_at = assignment.get("due_at")
                                    if (due_at and completion_map.get(assignment_id) and
                                        "✓ Submitted" not in existing_text):
                                        # Check if due date is within ±30 days
                                        try:
                                            due = datetime.fromisoformat(due_at.replace("Z", "+00:00"))
                                            days_diff = abs((due.date() - datetime.now(timezone.utc).date()).days)
                                            if days_diff <= 30:
                                                course_items.append(("mark_submitted", course_name, target_folder_id, name, existing_id))
                                            else:
                                                stats["skipped"] += 1
                                        except Exception:
                                            stats["skipped"] += 1
                                    else:
                                        stats["skipped"] += 1
                            else:
                                # No hash (legacy doc) -- still run the submitted check
                                assignment_id = str(assignment.get("id", ""))
                                due_at = assignment.get("due_at")
                                if (due_at and completion_map.get(assignment_id) and
                                    "✓ Submitted" not in existing_text):
                                    try:
                                        due = datetime.fromisoformat(due_at.replace("Z", "+00:00"))
                                        days_diff = abs((due.date() - datetime.now(timezone.utc).date()).days)
                                        if days_diff <= 30:
                                            course_items.append(("mark_submitted", course_name, target_folder_id, name, existing_id))
                                        else:
                                            stats["skipped"] += 1
                                    except Exception:
                                        stats["skipped"] += 1
                                else:
                                    stats["skipped"] += 1
                        except Exception as exc:
                            print(f"WARNING: could not check hash for '{name}': {exc}", file=sys.stderr)
                            stats["skipped"] += 1
                    continue

                # Locked assignments (confirmed live: locked_for_user=True) show no real
                # content via the API yet -- generating notes now would hit the same
                # "nothing to go on" case as an external-tool assignment, except here
                # it's guaranteed to have real content later. Just wait: no doc gets
                # created, so this is retried every run for free until Canvas reports
                # it unlocked. The calendar side (canvas-sync) is unaffected -- it syncs
                # from the ICS feed, which already includes locked assignments' due dates.
                if assignment.get("locked_for_user"):
                    stats["locked"] += 1
                    continue
                course_items.append(("assignment", course_name, target_folder_id, name, assignment))

            # Course files (PDF/PPTX) via module items -> Notes subfolder.
            # The /courses/{id}/files listing returns 403 for many courses (token
            # lacks the files scope or course has restricted file visibility), but
            # individual File items inside modules ARE accessible via content_id.
            # Scan every module's items and collect File entries.
            # Only scan modules if more than 6 hours have passed since last run
            last_modules_run = state.get("last_modules_run")
            should_scan_modules = True
            if last_modules_run:
                try:
                    last_run = datetime.fromisoformat(last_modules_run)
                    hours_since = (datetime.now(timezone.utc) - last_run).total_seconds() / 3600
                    should_scan_modules = hours_since >= 6
                except Exception:
                    should_scan_modules = True

            modules = []
            if should_scan_modules:
                try:
                    modules = canvas_get_all(f"/api/v1/courses/{course_id}/modules", {"per_page": 100})
                except AuthFailure:
                    raise
                except Exception as exc:
                    print(f"ERROR listing modules for course {course_id}: {exc}", file=sys.stderr)
                    stats["errors"] += 1
                    modules = []

            # Scan every module's items for File, Page, and ExternalUrl content.
            # (type, ref, title, module_name, module_folder_id): organize by module
            raw_module_items, seen_refs = [], set()
            module_folders = {}  # {module_name: folder_id}
            for mod in modules:
                if mod.get("items_count", 0) == 0:
                    continue
                module_name = mod.get("name", f"Module {mod.get('id')}")
                # Create module subfolder under Notes
                try:
                    module_folder_id = find_or_create_subfolder(all_docs, collection_id, notes_folder_id, module_name)
                    module_folders[module_name] = module_folder_id
                except Exception as exc:
                    print(f"ERROR creating module folder '{module_name}': {exc}", file=sys.stderr)
                    continue

                try:
                    items = canvas_get_all(
                        f"/api/v1/courses/{course_id}/modules/{mod['id']}/items",
                        {"per_page": 100},
                    )
                except AuthFailure:
                    raise
                except Exception:
                    continue
                for item in items:
                    itype = item.get("type")
                    if itype == "File":
                        ref = str(item.get("content_id") or "")
                    elif itype == "Page":
                        ref = item.get("page_url") or ""
                    elif itype == "ExternalUrl":
                        ref = item.get("external_url") or ""
                    else:
                        continue
                    key = (itype, ref)
                    if not ref or key in seen_refs:
                        continue
                    seen_refs.add(key)
                    raw_module_items.append((itype, ref, item.get("title") or ref, module_name, module_folder_id))

            # Collect all module folder IDs for searching
            all_module_parents = set(module_folders.values()) | {notes_folder_id}

            for itype, ref, title, module_name, module_folder_id in raw_module_items:
                if itype == "File":
                    try:
                        meta = canvas_get(f"/api/v1/files/{ref}")
                    except AuthFailure:
                        raise
                    except Exception:
                        continue
                    if meta.get("locked_for_user"):
                        continue
                    ct = meta.get("content-type") or meta.get("content_type") or ""
                    disp = meta.get("display_name") or title
                    name_lower = disp.lower()
                    if not (ct == "application/pdf" or name_lower.endswith(".pdf")
                            or "presentationml" in ct or name_lower.endswith(".pptx")):
                        continue
                    name = disp
                    # Check if item exists in any module folder or old flat Notes
                    existing_id, existing_parent = None, None
                    for d in all_docs:
                        if d.get("title") == name and d.get("parentDocumentId") in all_module_parents:
                            existing_id = d["id"]
                            existing_parent = d.get("parentDocumentId")
                            break
                    if existing_id is not None:
                        if existing_parent != module_folder_id:
                            # Wrong module folder - move it
                            course_items.append(("move", course_name, module_folder_id, name, existing_id))
                        else:
                            stats["skipped"] += 1
                        continue
                    course_items.append(("file", course_name, module_folder_id, name, meta))

                elif itype == "Page":
                    try:
                        page_data = canvas_get(f"/api/v1/courses/{course_id}/pages/{ref}")
                    except AuthFailure:
                        raise
                    except Exception:
                        continue
                    body_html = page_data.get("body") or ""
                    text = html_to_text(body_html)
                    if len(text) < 300:
                        continue  # skip near-empty pages (index, link-only, etc.)
                    name = page_data.get("title") or title
                    # Check if item exists in any module folder
                    existing_id, existing_parent = None, None
                    for d in all_docs:
                        if d.get("title") == name and d.get("parentDocumentId") in all_module_parents:
                            existing_id = d["id"]
                            existing_parent = d.get("parentDocumentId")
                            break
                    if existing_id is not None:
                        if existing_parent != module_folder_id:
                            course_items.append(("move", course_name, module_folder_id, name, existing_id))
                        else:
                            stats["skipped"] += 1
                        continue
                    course_items.append(("page", course_name, module_folder_id, name, text))

                elif itype == "ExternalUrl":
                    name = title
                    # Check if item exists in any module folder
                    existing_id, existing_parent = None, None
                    for d in all_docs:
                        if d.get("title") == name and d.get("parentDocumentId") in all_module_parents:
                            existing_id = d["id"]
                            existing_parent = d.get("parentDocumentId")
                            break
                    if existing_id is not None:
                        if existing_parent != module_folder_id:
                            course_items.append(("move", course_name, module_folder_id, name, existing_id))
                        else:
                            stats["skipped"] += 1
                        continue
                    yt_match = YT_URL_RE.search(ref)
                    if yt_match:
                        course_items.append(("youtube", course_name, module_folder_id, name, yt_match.group(1)))
                    else:
                        course_items.append(("url", course_name, module_folder_id, name, ref))

            if course_items:
                pending_by_course.append(course_items)

        # Round-robin across courses so every course gets fair progress
        round_robin = []
        while any(items for items in pending_by_course):
            for items in pending_by_course:
                if items:
                    round_robin.append(items.pop(0))

        for item_type, *item_args in round_robin:
            # Moves don't count toward MAX_NEW_PER_RUN -- they're cheap (no LLM call)
            if item_type == "move":
                course_name, parent_doc_id, name, item_data = item_args
                try:
                    outline_post("documents.move", {
                        "id": item_data,
                        "parentDocumentId": parent_doc_id,
                        "collectionId": collection_id,
                    })
                    for d in all_docs:
                        if d["id"] == item_data:
                            d["parentDocumentId"] = parent_doc_id
                            break
                    stats["moved"] += 1
                    print(f"moved '{course_name} / {name}' to correct location")
                except Exception as exc:
                    print(f"ERROR moving '{name}': {exc}", file=sys.stderr)
                    stats["errors"] += 1
                continue

            # mark_submitted is also cheap (no LLM call)
            if item_type == "mark_submitted":
                course_name, parent_doc_id, name, existing_id = item_args
                try:
                    doc_info = outline_post("documents.info", {"id": existing_id})
                    current_text = doc_info.get("text", "")
                    if "✓ Submitted" not in current_text:
                        new_text = current_text.replace("**Course:**", "**Status:** ✓ Submitted\n\n**Course:**", 1)
                        outline_post("documents.update", {"id": existing_id, "text": new_text, "publish": True, "icon": "✅"})
                        stats["marked_submitted"] += 1
                        print(f"marked '{course_name} / {name}' as submitted")
                except Exception as exc:
                    print(f"ERROR marking submitted '{name}': {exc}", file=sys.stderr)
                    stats["errors"] += 1
                continue

            if stats["created"] >= MAX_NEW_PER_RUN:
                stats["capped"] += 1
                continue

            try:
                if item_type == "assignment":
                    course_name, parent_doc_id, name, item_data = item_args
                    doc_text = generate_notes(course_name, item_data)
                    new_id = create_document(collection_id, parent_doc_id, name, doc_text, icon="📝")
                    all_docs.append({"id": new_id, "title": name, "parentDocumentId": parent_doc_id})
                    stats["created"] += 1
                    print(f"created assignment notes for '{course_name} / {name}'")

                elif item_type == "regen":
                    course_name, parent_doc_id, name, assignment, existing_id = item_args
                    generate_notes(course_name, assignment, existing_id=existing_id)
                    stats["regenerated"] += 1
                    print(f"regenerated notes for '{course_name} / {name}'")
                elif item_type == "file":
                    course_name, parent_doc_id, name, item_data = item_args
                    f = fetch_canvas_file(item_data["id"])
                    if f is None:
                        print(f"WARNING: skipping '{name}': file too large or unavailable", file=sys.stderr)
                        continue
                    doc_text = generate_presentation_notes(course_name, name, f["bytes"], f["content_type"])
                    if doc_text is None:
                        print(f"WARNING: skipping '{name}': no extractable content", file=sys.stderr)
                        continue
                    new_id = create_document(collection_id, parent_doc_id, name, doc_text)
                    all_docs.append({"id": new_id, "title": name, "parentDocumentId": parent_doc_id})
                    stats["created"] += 1
                    print(f"created file notes for '{course_name} / {name}'")

                elif item_type == "page":
                    course_name, parent_doc_id, name, item_data = item_args
                    # item_data is already the stripped page text
                    doc_text = generate_text_notes(course_name, name, item_data, "Canvas Page")
                    new_id = create_document(collection_id, parent_doc_id, name, doc_text)
                    all_docs.append({"id": new_id, "title": name, "parentDocumentId": parent_doc_id})
                    stats["created"] += 1
                    print(f"created page notes for '{course_name} / {name}'")

                elif item_type == "youtube":
                    course_name, parent_doc_id, name, item_data = item_args
                    # item_data is the YouTube video ID
                    transcript, reason = get_youtube_transcript(item_data)
                    if transcript is None:
                        print(f"WARNING: skipping YouTube '{name}': {reason}", file=sys.stderr)
                        continue
                    doc_text = generate_text_notes(course_name, name, transcript, "YouTube Video Transcript")
                    new_id = create_document(collection_id, parent_doc_id, name, doc_text)
                    all_docs.append({"id": new_id, "title": name, "parentDocumentId": parent_doc_id})
                    stats["created"] += 1
                    print(f"created YouTube notes for '{course_name} / {name}'")

                elif item_type == "url":
                    course_name, parent_doc_id, name, item_data = item_args
                    # item_data is the raw URL
                    text, reason = fetch_url_text(item_data)
                    if text is None:
                        print(f"WARNING: skipping URL '{name}': {reason}", file=sys.stderr)
                        continue
                    doc_text = generate_text_notes(course_name, name, text, "External Resource")
                    new_id = create_document(collection_id, parent_doc_id, name, doc_text)
                    all_docs.append({"id": new_id, "title": name, "parentDocumentId": parent_doc_id})
                    stats["created"] += 1
                    print(f"created URL notes for '{course_name} / {name}'")
            except AuthFailure:
                raise
            except Exception as exc:
                print(f"ERROR processing '{name}': {exc}", file=sys.stderr)
                stats["errors"] += 1
    except AuthFailure as exc:
        print(f"ERROR: {exc.service} appears invalid/expired: {exc.detail}", file=sys.stderr)
        raise_credential_alarm(exc.service, exc.detail)
        sys.exit(1)

    resolve_credential_alarm()

    # Update state timestamps after successful run
    now = datetime.now(timezone.utc).isoformat()
    state["last_assignments_run"] = now
    if should_scan_modules:
        state["last_modules_run"] = now
    write_state(state)

    print(
        f"canvas-outline-notes: created={stats['created']} moved={stats['moved']} "
        f"regenerated={stats['regenerated']} marked_submitted={stats['marked_submitted']} "
        f"skipped={stats['skipped']} locked-waiting={stats['locked']} "
        f"capped-for-next-run={stats['capped']} errors={stats['errors']}"
    )
    if stats["errors"]:
        sys.exit(1)


if __name__ == "__main__":
    main()
